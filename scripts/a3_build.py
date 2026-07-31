"""A3 report builder: JSON -> PPTX (single A3 landscape slide) and/or XLSX.

Usage:
    python a3_build.py a3.json --pptx out.pptx --xlsx out.xlsx
"""
import argparse
import json
import sys

ACCENT = (0xB9, 0x1C, 0x1C)
TYPE_LABELS = {
    "problem-solving": "问题解决型",
    "proposal": "提案型",
    "status": "状态汇报型",
}


def load(path):
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    for key in ("title", "sections"):
        if key not in data:
            sys.exit(f"ERROR: a3.json missing required key: {key}")
    return data


def meta_line(data):
    m = data.get("meta", {})
    parts = []
    for label, key in (("实践人", "author"), ("教练", "coach"),
                       ("部门", "department"), ("日期", "date"),
                       ("版本", "version")):
        if m.get(key):
            parts.append(f"{label}: {m[key]}")
    return "    ".join(parts)


def build_pptx(data, out_path):
    from pptx import Presentation
    from pptx.util import Mm, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    accent = RGBColor(*ACCENT)
    prs = Presentation()
    prs.slide_width = Mm(420)
    prs.slide_height = Mm(297)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Mm(12), Mm(8), Mm(300), Mm(14))
    tf = title_box.text_frame
    tf.text = f"{data['title']}  [{TYPE_LABELS.get(data.get('type'), 'A3')}]"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = accent

    meta_box = slide.shapes.add_textbox(Mm(240), Mm(10), Mm(168), Mm(10))
    mp = meta_box.text_frame.paragraphs[0]
    mp.text = meta_line(data)
    mp.font.size = Pt(10)
    mp.alignment = PP_ALIGN.RIGHT

    cols = {"left": [], "right": []}
    for s in data["sections"]:
        cols[s.get("side", "left")].append(s)

    top, bottom, gap = Mm(24), Mm(289), Mm(4)
    col_geom = {"left": (Mm(12), Mm(196)), "right": (Mm(212), Mm(196))}
    for side, sections in cols.items():
        if not sections:
            continue
        x, width = col_geom[side]
        avail = bottom - top - gap * (len(sections) - 1)
        box_h = int(avail / len(sections))
        y = top
        for s in sections:
            shape = slide.shapes.add_textbox(x, y, width, box_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
            shape.line.color.rgb = accent
            shape.line.width = Pt(1.25)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Mm(3)
            tf.margin_top = Mm(2)
            head = tf.paragraphs[0]
            head.text = s["title"]
            head.font.size = Pt(14)
            head.font.bold = True
            head.font.color.rgb = accent
            for line in s.get("content", "").split("\n"):
                para = tf.add_paragraph()
                para.text = line
                para.font.size = Pt(10.5)
            y += box_h + gap

    prs.save(out_path)
    print(f"OK: PPTX written -> {out_path}")


def build_xlsx(data, out_path):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    accent_hex = "B91C1C"
    wb = Workbook()
    ws = wb.active
    ws.title = "A3报告"
    thin = Side(style="thin", color="D4D4D4")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    wrap = Alignment(wrap_text=True, vertical="top")

    for col, width in (("A", 55), ("B", 3), ("C", 55)):
        ws.column_dimensions[col].width = width

    ws.merge_cells("A1:C1")
    c = ws["A1"]
    c.value = f"{data['title']}  [{TYPE_LABELS.get(data.get('type'), 'A3')}]"
    c.font = Font(size=16, bold=True, color="FFFFFF")
    c.fill = PatternFill("solid", fgColor=accent_hex)
    c.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 30

    ws.merge_cells("A2:C2")
    ws["A2"].value = meta_line(data)
    ws["A2"].alignment = Alignment(horizontal="right")
    ws["A2"].font = Font(size=9, color="555555")

    cols = {"left": [], "right": []}
    for s in data["sections"]:
        cols[s.get("side", "left")].append(s)

    def write_col(letter, sections):
        row = 4
        for s in sections:
            head = ws[f"{letter}{row}"]
            head.value = s["title"]
            head.font = Font(bold=True, size=12, color=accent_hex)
            head.border = border
            body = ws[f"{letter}{row + 1}"]
            body.value = s.get("content", "")
            body.alignment = wrap
            body.border = border
            lines = s.get("content", "").count("\n") + 1
            ws.row_dimensions[row + 1].height = max(
                ws.row_dimensions[row + 1].height or 0, lines * 15 + 6)
            row += 3

    write_col("A", cols["left"])
    write_col("C", cols["right"])

    ws.print_area = f"A1:C{ws.max_row}"
    ws.page_setup.orientation = "landscape"
    ws.page_setup.paperSize = 8  # A3
    wb.save(out_path)
    print(f"OK: XLSX written -> {out_path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("json_file")
    ap.add_argument("--pptx")
    ap.add_argument("--xlsx")
    args = ap.parse_args()
    if not args.pptx and not args.xlsx:
        sys.exit("ERROR: pass --pptx and/or --xlsx output path")
    data = load(args.json_file)
    if args.pptx:
        build_pptx(data, args.pptx)
    if args.xlsx:
        build_xlsx(data, args.xlsx)


if __name__ == "__main__":
    main()
